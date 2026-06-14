import os
import json
from pptx import Presentation
from pptx.util import Pt
import io
import base64
import tempfile
import shutil
from typing import List
import fitz  # PyMuPDF
from fastapi import UploadFile
from langfuse.openai import AsyncOpenAI  # Restaurado conforme pedido
from langfuse import get_client
from langchain_postgres.vectorstores import PGVector
from langchain_openai import OpenAIEmbeddings
from dotenv import load_dotenv
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Pt, Inches

load_dotenv()

# Cliente OpenAI padrão com timeout aumentado
client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"), timeout=120.0)

# Inicialização Lazy do PGVector (para não travar a inicialização do app)
_db_bncc = None
_db_aee = None
_embeddings = None

def get_embeddings():
    global _embeddings
    if _embeddings is None:
        print("🔍 [LOG] Inicializando OpenAI Embeddings...")
        _embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
    return _embeddings

def get_vectorstore(collection_name: str) -> PGVector:
    print(f"🔍 [LOG] Conectando ao PGVector na coleção '{collection_name}'...")
    conn_string = PGVector.connection_string_from_db_params(
        host=os.getenv("VECTOR_HOST"),
        port=os.getenv("VECTOR_PORT"),
        database=os.getenv("VECTOR_NAME"),
        user=os.getenv("VECTOR_USER"),
        password=os.getenv("VECTOR_PASSWORD"),
        driver="psycopg",
    )
    
    return PGVector(
        connection=conn_string, 
        collection_name=collection_name, 
        embeddings=get_embeddings(), 
        create_extension=False
    )

async def convert_files_to_base64_images(files: List[UploadFile]) -> List[str]:
    """Transforma tudo (inclusive PDFs) em imagens puras para o Vision do GPT-4o."""
    arquivos_validos = [f for f in files if f.filename and f.filename.strip() != ""]
    if not arquivos_validos:
        return []

    base64_images = []
    temp_dir = tempfile.mkdtemp()
    
    try:
        for file in arquivos_validos[:5]: # Máximo de 5 arquivos
            file_path = os.path.join(temp_dir, file.filename)
            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)
            
            if file.filename.lower().endswith(".pdf"):
                doc = fitz.open(file_path)
                # Limita a 3 páginas por PDF e reduz resolução para evitar timeout
                for page_num in range(min(len(doc), 3)):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5)) # Reduzido de 2,2 para 1,5,1,5
                    img_bytes = pix.tobytes("jpeg", jpg_quality=70) # Qualidade reduzida para 70
                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                    base64_images.append(img_b64)
                doc.close() 
            else:
                file.file.seek(0)
                img_bytes = file.file.read()
                img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                base64_images.append(img_b64)
    except Exception as e:
        print(f"❌ [ERRO NO PROCESSAMENTO DAS IMAGENS]: {e}")
        raise e
    finally:
        shutil.rmtree(temp_dir)
        
    return base64_images

async def processar_e_gerar_aula(tipo: str, materia: str, ano: str, necessidades: str, prompt_chat: str, arquivos: List[UploadFile]) -> str:
    """Função mestre chamada pela rota /gerar-aula do main.py"""
    global _db_bncc, _db_aee
    
    print(f"🚀 [LOG] INICIANDO GERAÇÃO: {materia} | {ano} | AEE: {necessidades}")
    
    try:
        # 👇 A SOLUÇÃO DO OPERATIONAL ERROR: Criamos a conexão na hora do uso, sem variáveis 'global'
        db_bncc = get_vectorstore("bncc_embeddings_openai")
        
        # 1. Busca RAG
        query_rag = f"{materia} {ano}: {prompt_chat}"
        docs_bncc = db_bncc.similarity_search(query_rag, k=3)
        contexto_bncc = "\n".join([d.page_content for d in docs_bncc])
        
        contexto_aee = ""
        if necessidades and necessidades.strip() != "":
            db_aee = get_vectorstore("diretrizes_aee_openai")
            docs_aee = db_aee.similarity_search(f"Adaptação AEE: {necessidades}", k=3)
            contexto_aee = "\n".join([d.page_content for d in docs_aee])

        # 2. Transforma arquivos anexados em imagens puras
        imagens_b64 = await convert_files_to_base64_images(arquivos)

        # 3. Base do Contexto (AGORA COM AEE NO NÚCLEO)
        base_context = f"""
        DISCIPLINA: {materia} | PÚBLICO: {ano}
        ALUNOS DE INCLUSÃO (AEE): {necessidades if necessidades.strip() != "" else "Turma Regular (Sem AEE)"}
        DIRETRIZES BNCC: {contexto_bncc}
        DIRETRIZES AEE: {contexto_aee if contexto_aee else "Padrão."}
        """

        # --- REGRAS DE OURO (IDIOMA E DIFICULDADE) ---
        regra_idioma = ""
        if materia.lower() in ["inglês", "língua inglesa", "english"]:
            regra_idioma = "REGRA DE OURO: Como a matéria é Inglês, TODO O CONTEÚDO (explicações, enunciados, questões, objetivos) DEVE SER EM INGLÊS. Use um vocabulário rico e adequado ao nível do aluno."
        else:
            regra_idioma = "O conteúdo deve ser em Português do Brasil."

        regra_dificuldade = f"As atividades e avaliações devem ter um nível de profundidade e desafio elevado para o ano {ano}, incentivando o pensamento crítico e a resolução de problemas complexos, sem subestimar a capacidade do aluno, mas respeitando a BNCC."

        print("🧠 [LOG] Montando Super Prompt Visual e enviando para gpt-4o...")
        print(base_context)
        
        # --- 1. PLANO DE AULA ---
        if tipo == "Aula":
            system_instruction = f"""
            Você é um Coordenador Pedagógico Master Class focado em criar roteiros EXAUSTIVOS.
            {base_context}
            
            {regra_idioma}
            {regra_dificuldade}

            REGRAS DE OURO DA EXCELÊNCIA (OBRIGATÓRIO):
            1. 🎯 OBJETIVOS DA BNCC: Logo no início do documento, ANTES de começar o plano de aula, crie uma seção "Objetivos da BNCC" e liste os códigos alfanuméricos e habilidades trabalhadas com base nas diretrizes fornecidas no contexto.
            2. 👁️ LEITURA DE IMAGENS: Analise as fotos anexadas e utilize o conteúdo delas como base principal.
            3. 🧠 SEQUÊNCIA: Aquecimento -> Explicação -> Prática Guiada -> Fechamento.
            4. 📖 DETALHE ABSURDO: Descreva o passo a passo (Ex: escreva na lousa).
            5. 🔀 VARIEDADE: Alterne dinâmicas diariamente.
            6. 🎵 MÍDIA: Sugira vídeos/músicas com termos de busca reais.
            7. 🛠️ MATERIAIS: Comandos exatos na língua da matéria.
            8. ⏱️ TEMPO: Aloque o tempo realista (45-50 min).
            9. 🤐 MODO DOCUMENTO: Apenas o documento final, sem conversas.
            """
            
            if necessidades and necessidades.strip() != "":
                system_instruction += """
            10. 🚨 REGRA ABSOLUTA DE AEE: Você DEVE OBRIGATORIAMENTE criar uma seção ao final chamada "--- ADAPTAÇÕES PARA AEE ---" detalhando as adaptações para os alunos citados no contexto.
                """

        # --- 2. ATIVIDADE DE FIXAÇÃO ---
        elif tipo == "Atividade":
            system_instruction = f"""
            Você é um Criador de Avaliações e Exercícios Escolares de Elite.
            {base_context}
            
            {regra_idioma}
            {regra_dificuldade}

            🚨 DIRETRIZ SUPREMA (PEDIDO DO PROFESSOR):
            "{prompt_chat}"
            ATENÇÃO: A instrução acima é a LEI MÁXIMA deste documento. Se o professor pedir um número específico de questões ou um tipo específico, OBEDEÇA. Esta ordem anula as regras gerais abaixo em caso de conflito.
            
            REGRAS GERAIS:
            1. IMAGENS E ILUSTRAÇÕES: VOCÊ NÃO PODE GERAR IMAGENS REAIS. Sempre que precisar de uma imagem, crie "Caixas de Desenho" para o aluno interagir (Ex: "[ Desenhe um Robot aqui ]") OU deixe um espaço demarcado `[ PROFESSOR: COLE A IMAGEM AQUI ]`.
            2. VARIEDADE: Mescle os tipos de questões (múltipla escolha, verdadeiro/falso, completar e dissertativas).
            3. GABARITO: Forneça o gabarito exclusivo do professor no final do documento.
            """

        # --- 3. TAREFA DE CASA ---
        elif tipo == "Tarefas":
            system_instruction = f"""
            Você é um Pedagogo especialista em lições de casa eficientes e desafiadoras.
            {base_context}
            
            {regra_idioma}
            {regra_dificuldade}

            OBJETIVO: Criar uma Tarefa de Casa envolvente e de curta duração baseada no conteúdo das imagens ou contexto.
            REGRAS:
            1. INSTRUÇÕES CLARAS: Crie um cabeçalho para o aluno. As instruções devem ser simples e diretas.
            2. TEMPO: A tarefa deve levar no máximo 20 minutos para ser feita em casa.
            3. ENVOLVIMENTO FAMILIAR: Se possível para a idade, inclua uma pequena etapa onde o aluno precisa perguntar algo aos pais.
            4. MODO DOCUMENTO: Entregue a folha pronta para impressão ou cópia na agenda. Sem conversas paralelas.
            """

        # --- 4. PROVA / AVALIAÇÃO ---
        elif tipo == "Prova":
            system_instruction = f"""
            Você é um Especialista em Avaliação Escolar de Alto Nível (Banca Examinadora).
            {base_context}
            
            {regra_idioma}
            {regra_dificuldade}

            🚨 DIRETRIZ SUPREMA (PEDIDO DO PROFESSOR):
            "{prompt_chat}"
            ATENÇÃO: A instrução acima é a LEI MÁXIMA deste documento. 
            
            PRIORIDADE ABSOLUTA: Se o professor enviou imagens ou arquivos, considere que aquilo é uma PROVA que ele já usa. Sua missão é ADAPTAR essa prova para o nível de desafio exigido ({ano}) e também criar as adaptações de AEE necessárias.
            
            REGRAS GERAIS:
            1. ESTRUTURA FORMAL: Inclua cabeçalho completo (Escola, Nome, Turma, Data, Nota).
            2. IMAGENS NA PROVA: Deixe áreas demarcadas `[ PROFESSOR: INSERIR IMAGEM AQUI ]`.
            3. QUESTÕES PADRÃO: Se o professor NÃO especificou formato, crie 7 questões (5 objetivas e 2 discursivas) com alto nível de exigência.
            4. GABARITO: Entregue o gabarito completo no final do documento.
            """

        # --- 5. TRABALHO INDIVIDUAL ---
        elif tipo == "Individual":
            system_instruction = f"""
            Você é um Professor inovador criando um Projeto Individual de alto nível.
            {base_context}
            
            {regra_idioma}
            {regra_dificuldade}

            OBJETIVO: Criar um trabalho de pesquisa ou projeto "Mão na Massa" desafiador que o aluno fará sozinho.
            REGRAS:
            1. ESCOPO: Defina o tema central baseado no conteúdo anexado ou contexto.
            2. PASSO A PASSO PARA O ALUNO: Descreva exatamente o que o aluno deve pesquisar, construir ou escrever com profundidade.
            3. FORMATO DE ENTREGA: Defina como será avaliado (cartolina, maquete, redação, apresentação).
            4. RUBRICA DE AVALIAÇÃO: Crie uma tabela textual com critérios claros e exigentes.
            """

        # --- 6. TRABALHO EM GRUPO ---
        elif tipo == "Grupo":
            system_instruction = f"""
            Você é um Especialista em Metodologias Ativas e Aprendizagem Colaborativa.
            {base_context}
            
            {regra_idioma}
            {regra_dificuldade}

            OBJETIVO: Estruturar um Trabalho em Grupo complexo e estimulante.
            REGRAS:
            1. DIVISÃO DE PAPÉIS: Especifique papéis (Ex: O Pesquisador, O Redator, O Apresentador, O Designer).
            2. OBJETIVO DO GRUPO: Qual é o problema complexo que eles devem resolver?
            3. CRONOGRAMA: Divida o tempo em sala ou prazos de entrega.
            4. RUBRICA: Como o grupo será avaliado de forma rigorosa.
            """

        # --- 7. PAUTA DE REUNIÃO ---
        elif tipo == "Pauta":
            system_instruction = f"""
            Você é um Gestor Escolar estruturando uma reunião oficial.
            {base_context}
            
            OBJETIVO: Criar uma Pauta de Reunião (com Pais ou Coordenação) estruturada e profissional.
            REGRAS:
            1. ESTRUTURA: Inclua Data, Local, Participantes, Objetivo da Reunião.
            2. TÓPICOS: Liste os tópicos que serão abordados em formato de "Bullet Points" com estimativa de tempo (Ex: "Desempenho da turma no semestre - 15 min").
            3. REGISTROS: Deixe um bloco pontilhado ou vazio para "Anotações e Deliberações".
            4. MODO DOCUMENTO: Tom sério, empático e diretivo.
            """

        # --- 8. PARECER PEDAGÓGICO ---
        elif tipo == "Parecer":
            system_instruction = f"""
            Você é um Orientador Educacional elaborando um documento de altíssima responsabilidade.
            {base_context}
            
            OBJETIVO: Escrever um Parecer Pedagógico Descritivo sobre um aluno ou turma.
            REGRAS:
            1. LINGUAGEM TÉCNICA E AFETIVA: Use vocabulário pedagógico formal, mas sem perder o acolhimento. Evite rótulos negativos (em vez de "o aluno é agressivo", use "o aluno encontra desafios na autorregulação emocional").
            2. SEÇÕES DO DOCUMENTO: 
               - Desenvolvimento Cognitivo e Acadêmico.
               - Interação Social e Comportamento.
               - Adaptações (Foque pesado nas DIRETRIZES INCLUSÃO/AEE se houver, relatando como estão sendo aplicadas).
               - Recomendações para o próximo ciclo/família.
            3. MODO DOCUMENTO: Entregue pronto para ser impresso, assinado e entregue à diretoria/pais.
            """

        # --- FAILSAFE (Caso algum tipo passe em branco) ---
        else:
            system_instruction = f"""
            Você é um Assistente Educacional Pleno. Crie o documento: {tipo}.
            {base_context}
            Baseie TODA a atividade nas imagens fornecidas. Entregue um material denso, livre de conversas iniciais ou finais.
            """

        # 4. Montagem da Mensagem com Visão Ocular
        conteudo_usuario = [{"type": "text", "text": f"Pedido do professor: {prompt_chat}\nAnalise as páginas anexadas e gere o material exigido."}]
        
        for b64 in imagens_b64:
            conteudo_usuario.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{b64}"}
            })
        
        # Injeção da Regra de Ouro do AEE
        if necessidades and necessidades.strip() != "":
            if tipo not in ["Pauta", "Parecer"]: # Não precisa de atividade adaptada em pauta de reunião
                system_instruction += f"\n\n🚨 REGRA VITAL DE INCLUSÃO: A turma possui alunos com {necessidades}. Você DEVE OBRIGATORIAMENTE criar uma seção no final do documento chamada '--- VERSÃO ADAPTADA (AEE) ---' com o conteúdo redesenhado, simplificado e visualmente adaptado para eles."

        full_messages = [
            {"role": "system", "content": system_instruction},
            {"role": "user", "content": conteudo_usuario}
        ]
        
        response = await client.chat.completions.create(
            model="gpt-4o",
            messages=full_messages,
            temperature=0
        )

        try:
            # flush é importante para garantir que o trace chegue ao Langfuse
            get_client().flush()
            print("✅ [LOG] Trace enviado ao Langfuse com sucesso!")
        except Exception as e:
            print(f"⚠️ [AVISO]: Falha ao enviar para Langfuse, mas a aula foi gerada: {e}")

        print("🎉 [LOG] GERAÇÃO CONCLUÍDA COM SUCESSO!")
        return response.choices[0].message.content

    except Exception as e:
        print(f"❌ [ERRO CRÍTICO NA GERAÇÃO]: {str(e)}")
        raise e

async def processar_cascata(tipo: str, conteudo_base: str) -> str:
    """Motor que gera materiais baseados em um plano, com foco em AEE e Design."""
    print(f"🚀 [LOG] GERAÇÃO EM CASCATA: {tipo}")
    
    if tipo == "SlideText":
        system_instruction = f"""
        Você é um Designer Pedagógico de Elite e Especialista em Comunicação Visual.
        Sua missão é transformar o PLANO DE AULA base em uma apresentação de PowerPoint de ALTO IMPACTO, digna de uma conferência educacional.
        
        🚨 REGRA DE OURO (IDIOMA): 
        O material gerado DEVE seguir o idioma do PLANO BASE. Se a aula é de English (Inglês), os slides devem ser em Inglês.
        
        DIRETRIZES DE OURO:
        1. MINIMALISMO E IMPACTO: Menos texto, mais significado. Use frases curtas que instiguem o aluno.
        2. HIERARQUIA VISUAL: O título deve ser provocador. Os tópicos devem ser a essência do aprendizado.
        3. STORYTELLING PEDAGÓGICO:
           - Slide 1: Capa Impactante.
           - Slide 2: "O que vamos descobrir hoje?" (Objetivos).
           - Slides de Desenvolvimento: Conteúdo denso mas fatiado.
           - Slides AEE: Suporte visual simplificado, focado em ícones mentais e associações claras.
           - Slide Final: Desafio do Dia ou "Momento Reflexão".
        
        ESTRUTURA DE LAYOUT (layout):
        - "capa": Título Grande e Subtítulo.
        - "conteudo": Título e até 4 tópicos curtos.
        - "aee": Slides adaptados com linguagem ultra-simples e destaque para conceitos-chave.
        - "pergunta": Slide com uma grande questão central para debate.
        
        FORMATO DE SAÍDA:
        Retorne APENAS um Array JSON puro (sem markdown, sem ```json) com esta estrutura:
        [
          {{"titulo": "TÍTULO PROVOCADOR", "subtitulo": "Disciplina e Série", "layout": "capa"}},
          {{"titulo": "O DESAFIO DE HOJE", "topicos": ["Ponto 1", "Ponto 2"], "layout": "conteudo"}},
          {{"titulo": "EM RESUMO (AEE)", "topicos": ["Ideia Central 1", "Associação 2"], "layout": "aee"}}
        ]
        """
    elif tipo in ["Prova", "Tarefas", "Atividade", "Individual", "Grupo", "Trabalho"]:
        system_instruction = f"""
        [COMANDO DE ALTA PRIORIDADE - MODO FOLHA DE IMPRESSÃO]
        Você receberá um PLANO DE AULA longo. 
        
        🚨 REGRA DE OURO (IDIOMA): 
        O material gerado DEVE seguir o idioma do PLANO BASE. Se a aula é de English (Inglês), o material deve ser em Inglês.
        
        🚨 REGRA DE DIFICULDADE: 
        Mantenha um alto nível de desafio e profundidade pedagógica, adequado ao ano escolar citado no plano, incentivando o raciocínio crítico.

        AÇÃO PROIBIDA: VOCÊ ESTÁ ESTRITAMENTE PROIBIDO DE GERAR "AULAS" (Aula 1, Aula 2, etc). Não cite roteiros, não cite oração, não cite cronômetro.
        AÇÃO OBRIGATÓRIA: Crie UMA ÚNICA FOLHA PARA O ALUNO do tipo: {tipo}.
        DESIGN OBRIGATÓRIO: Use Markdown avançado. Crie tabelas, listas, linhas em branco para o aluno escrever (______) e caixas de seleção [ ].
        """
        
        if tipo == "Prova":
            system_instruction += "\nESTRUTURA: Cabeçalho (Nome/Data/Nota). 5 questões objetivas e 2 discursivas. No final, Gabarito e uma versão adaptada para AEE (Inclusão)."
        elif tipo == "Tarefas":
            system_instruction += "\nESTRUTURA: Lição de Casa de 1 página. Exercícios diretos (ligar, completar, desenhar)."
        elif tipo == "Trabalho":
            system_instruction += "\nESTRUTURA: Trabalho de Pesquisa ou Prático. Defina Tema, Instruções, O que entregar e Rubrica de Avaliação."
        
        system_instruction += "\n\n🚨 REGRA VITAL DE INCLUSÃO: Analise o PLANO BASE. Se houver menção a alunos de inclusão/AEE, você DEVE OBRIGATORIAMENTE criar uma seção no final do documento chamada '--- VERSÃO ADAPTADA (AEE) ---' com a atividade simplificada."
    else:
        system_instruction = f"Você é um assistente pedagógico. Gere um conteúdo do tipo {tipo} baseado no plano fornecido."
            
    messages = [
        {"role": "system", "content": system_instruction},
        {"role": "user", "content": f"PLANO BASE:\n\n{conteudo_base}"}
    ]
    
    response = await client.chat.completions.create(model="gpt-4o", messages=messages, temperature=0)
    
    try:
        get_client().flush()
    except Exception as e:
        print(f"⚠️ [AVISO]: Falha ao enviar para Langfuse, mas a cascata foi gerada: {e}")

    return response.choices[0].message.content

async def gerar_arquivo_pptx(conteudo_base: str, logo_b64: str = None) -> str:
    """Gera um PowerPoint de alto nível, com design limpo e tipografia impactante."""
    print("🚀 [LOG] GERANDO APRESENTAÇÃO DE ELITE (.pptx)...")
    
    try:
        json_text = await processar_cascata("SlideText", conteudo_base)
    except Exception as e:
        print(f"❌ [ERRO AO CHAMAR IA PARA SLIDES]: {e}")
        raise e
    
    try:
        # Extração Robusta de JSON
        if "[" in json_text and "]" in json_text:
            json_text = json_text[json_text.find("["):json_text.rfind("]")+1]
            
        json_text = json_text.replace("```json", "").replace("```", "").strip()
        slides_data = json.loads(json_text)
        print(f"✅ [LOG] {len(slides_data)} slides preparados para montagem.")
    except Exception as e:
        print(f"❌ [ERRO JSON SLIDES]: {json_text}")
        raise Exception("Erro ao estruturar slides na IA. Tente novamente.")

    # 1. Carrega o Template Profissional
    template_path = os.path.join(os.path.dirname(__file__), "template_lecionia.pptx")
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"⚠️ [AVISO]: Template não encontrado. Usando layout padrão. Erro: {e}")
        prs = Presentation() 

    # 2. Prepara a Logo da Instituição
    logo_stream = None
    if logo_b64 and "base64," in logo_b64:
        try:
            base64_data = logo_b64.split(",")[1]
            logo_bytes = base64.b64decode(base64_data)
            logo_stream = io.BytesIO(logo_bytes)
        except:
            pass

    # 3. Montagem Cirúrgica dos Slides
    for data in slides_data:
        l_type = data.get("layout", "conteudo").lower()
        
        # Mapeamento Inteligente de Layouts
        # 0: Capa | 1: Título e Conteúdo | 2: Layout Alternativo (AEE) | 3: Pergunta (se houver)
        if l_type == "capa":
            layout_idx = 0
        elif l_type == "aee":
            layout_idx = 2 if len(prs.slide_layouts) > 2 else 1
        elif l_type == "pergunta":
            layout_idx = 3 if len(prs.slide_layouts) > 3 else 1
        else:
            layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
        
        try:
            slide_layout = prs.slide_layouts[layout_idx]
        except:
            slide_layout = prs.slide_layouts[0] 
            
        slide = prs.slides.add_slide(slide_layout)
        
        # --- ESTILIZAÇÃO DO TÍTULO ---
        if slide.shapes.title:
            slide.shapes.title.text = data.get("titulo", "").upper()
            tf = slide.shapes.title.text_frame
            for paragraph in tf.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if layout_idx in [0, 3] else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.name = 'Calibri'
                    run.font.size = Pt(54) if layout_idx == 0 else Pt(44)
                    if l_type == "pergunta":
                        run.font.color.rgb = RGBColor(124, 58, 237) # Roxo LecionaAI
        
        # --- ESTILIZAÇÃO DO SUBTÍTULO (Capa) ---
        if layout_idx == 0 and "subtitulo" in data:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 4: # SUBTITLE
                    shape.text = data["subtitulo"]
                    tf_sub = shape.text_frame
                    tf_sub.paragraphs[0].alignment = PP_ALIGN.CENTER
                    for run in tf_sub.paragraphs[0].runs:
                        run.font.size = Pt(24)
                        run.font.color.rgb = RGBColor(100, 116, 139)
                    break

        # --- ESTILIZAÇÃO DOS TÓPICOS ---
        if "topicos" in data and len(data["topicos"]) > 0:
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type in [2, 7]: # BODY ou OBJECT
                    body_shape = shape
                    break
            
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                
                for topico in data["topicos"]:
                    p = tf.add_paragraph()
                    p.text = f"  {topico}"
                    p.space_before = Pt(12)
                    p.level = 0
                    
                    # Estilo Visual por Tipo
                    if l_type == "aee":
                        p.font.size = Pt(32)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0, 102, 204)
                    else:
                        p.font.size = Pt(28)
                        p.font.color.rgb = RGBColor(31, 41, 55)

        # --- POSICIONAMENTO DA LOGO (Branding) ---
        if logo_stream:
            logo_stream.seek(0)
            try:
                # Top-Right corner with safe margins
                height = Inches(0.7)
                left = prs.slide_width - Inches(1.8)
                top = Inches(0.3)
                slide.shapes.add_picture(logo_stream, left, top, height=height)
            except:
                pass

    # 4. Finalização do Stream
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    
    return base64.b64encode(pptx_stream.read()).decode('utf-8')